"""Generate Claude Code Workshop PowerPoint presentation."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.oxml.ns import qn
from pptx.enum.dml import MSO_THEME_COLOR
import copy
from lxml import etree

# ── Colour palette ────────────────────────────────────────────────────────────
BG       = RGBColor(0x0D, 0x11, 0x17)   # near-black background
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)   # title text
GREY     = RGBColor(0xC9, 0xD1, 0xD9)  # body text
ORANGE   = RGBColor(0xE9, 0x7B, 0x3C)  # accent / Anthropic brand
TEAL     = RGBColor(0x2D, 0xD4, 0xBF)  # secondary accent
CODE_BG  = RGBColor(0x16, 0x1B, 0x22)  # code-block background
HL_BG    = RGBColor(0x1C, 0x2B, 0x3A)  # highlight box background
DIM_GREY = RGBColor(0x58, 0x65, 0x69)  # table border / dim text
DARK_ROW = RGBColor(0x13, 0x1A, 0x21)  # alternate table row

# ── Slide dimensions (16:9 widescreen) ───────────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)

MARGIN = Inches(0.55)
TITLE_Y = Inches(0.38)
TITLE_H = Inches(0.75)
CONTENT_Y = Inches(1.25)
CONTENT_H = Inches(5.9)


# ── Helpers ───────────────────────────────────────────────────────────────────

def new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


def blank_slide(prs: Presentation):
    blank_layout = prs.slide_layouts[6]  # completely blank
    return prs.slides.add_slide(blank_layout)


def add_background(slide):
    """Fill whole slide with near-black."""
    bg = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        0, 0, W, H
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG
    bg.line.fill.background()  # no border
    # push to back
    slide.shapes._spTree.remove(bg._element)
    slide.shapes._spTree.insert(2, bg._element)


def add_textbox(slide, text: str, x, y, w, h,
                font_size=18, bold=False, color=WHITE,
                align=PP_ALIGN.LEFT, font_name="Calibri",
                word_wrap=True) -> object:
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    return txBox


def add_title(slide, title_text: str, font_size=32):
    tb = add_textbox(
        slide, title_text,
        MARGIN, TITLE_Y,
        W - 2 * MARGIN, TITLE_H,
        font_size=font_size, bold=True,
        color=WHITE, font_name="Calibri"
    )
    return tb


def add_orange_rule(slide, y=Inches(1.15), width=None):
    """Horizontal orange accent line below title."""
    rule_w = width or (W - 2 * MARGIN)
    rule = slide.shapes.add_shape(1, MARGIN, y, rule_w, Pt(2))
    rule.fill.solid()
    rule.fill.fore_color.rgb = ORANGE
    rule.line.fill.background()


def add_bullet_box(slide, bullets: list[str], x, y, w, h,
                   font_size=15, color=GREY, bullet_char="•"):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for bullet in bullets:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = f"{bullet_char}  {bullet}"
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.name = "Calibri"
    return txBox


def add_code_block(slide, code_text: str, x, y, w, h, font_size=10):
    """Dark rounded rectangle with monospace code text."""
    shape = slide.shapes.add_shape(
        5,  # rounded rectangle
        x, y, w, h
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = CODE_BG
    shape.line.color.rgb = RGBColor(0x30, 0x38, 0x42)
    shape.line.width = Pt(0.75)

    tf = shape.text_frame
    tf.word_wrap = True
    # set margins
    tf.margin_left   = Inches(0.15)
    tf.margin_right  = Inches(0.15)
    tf.margin_top    = Inches(0.12)
    tf.margin_bottom = Inches(0.12)

    lines = code_text.strip().split("\n")
    first = True
    for line in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.name = "Courier New"
        run.font.color.rgb = RGBColor(0xA5, 0xD6, 0xFF)
    return shape


def add_highlight_box(slide, text: str, x, y, w, h,
                      font_size=13, text_color=WHITE):
    """Blue-dark box with orange left border stripe."""
    # Main box
    box = slide.shapes.add_shape(1, x, y, w, h)
    box.fill.solid()
    box.fill.fore_color.rgb = HL_BG
    box.line.color.rgb = ORANGE
    box.line.width = Pt(1)

    # Orange left stripe
    stripe = slide.shapes.add_shape(1, x, y, Inches(0.08), h)
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = ORANGE
    stripe.line.fill.background()

    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.18)
    tf.margin_right = Inches(0.12)
    tf.margin_top = Inches(0.1)
    tf.margin_bottom = Inches(0.1)

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = text_color
    run.font.name = "Calibri"
    run.font.italic = True
    return box


def set_cell_bg(cell, rgb: RGBColor):
    """Set table cell background colour via XML."""
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    solidFill = etree.SubElement(tcPr, qn("a:solidFill"))
    srgbClr = etree.SubElement(solidFill, qn("a:srgbClr"))
    # str(RGBColor) returns the 6-char hex value in this pptx version
    srgbClr.set("val", str(rgb).upper())


def set_cell_text(cell, text: str, font_size=12, bold=False,
                  color=GREY, font_name="Calibri", align=PP_ALIGN.LEFT):
    cell.text = text
    p = cell.text_frame.paragraphs[0]
    p.alignment = align
    if p.runs:
        run = p.runs[0]
    else:
        run = p.add_run()
        run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    cell.text_frame.word_wrap = True


def add_table(slide, headers: list[str], rows: list[list[str]],
              x, y, w, h, col_widths=None):
    num_cols = len(headers)
    num_rows = len(rows) + 1  # +1 for header
    tbl = slide.shapes.add_table(num_rows, num_cols, x, y, w, h).table

    # Set column widths
    if col_widths:
        for i, cw in enumerate(col_widths):
            tbl.columns[i].width = cw

    # Header row
    for ci, hdr in enumerate(headers):
        cell = tbl.cell(0, ci)
        set_cell_bg(cell, RGBColor(0x1E, 0x2C, 0x3D))
        set_cell_text(cell, hdr, font_size=12, bold=True, color=ORANGE)

    # Data rows
    for ri, row in enumerate(rows):
        bg = DARK_ROW if ri % 2 == 0 else RGBColor(0x16, 0x1C, 0x27)
        for ci, val in enumerate(row):
            cell = tbl.cell(ri + 1, ci)
            set_cell_bg(cell, bg)
            set_cell_text(cell, val, font_size=11, color=GREY)

    return tbl


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE BUILDERS
# ══════════════════════════════════════════════════════════════════════════════

def slide_01(prs):
    """Title slide."""
    slide = blank_slide(prs)
    add_background(slide)

    # Large title
    title_tb = slide.shapes.add_textbox(MARGIN, Inches(1.8), W - 2 * MARGIN, Inches(1.2))
    tf = title_tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "Claude Code Workshop"
    run.font.size = Pt(52)
    run.font.bold = True
    run.font.color.rgb = WHITE
    run.font.name = "Calibri"

    # Orange horizontal rule under title
    add_orange_rule(slide, y=Inches(3.1), width=Inches(8))
    # Center the rule
    for shp in slide.shapes:
        if shp.shape_type == 1 and shp.fill.fore_color.rgb == ORANGE:
            shp.left = int((W - Inches(8)) / 2)

    # Subtitle
    sub_tb = slide.shapes.add_textbox(MARGIN, Inches(3.25), W - 2 * MARGIN, Inches(0.8))
    tf2 = sub_tb.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = "Building AI-Powered Applications with Agents, Skills & MCP"
    run2.font.size = Pt(20)
    run2.font.color.rgb = GREY
    run2.font.name = "Calibri"

    # Bottom tag line
    tag_tb = slide.shapes.add_textbox(MARGIN, Inches(6.6), W - 2 * MARGIN, Inches(0.5))
    tf3 = tag_tb.text_frame
    p3 = tf3.paragraphs[0]
    p3.alignment = PP_ALIGN.CENTER
    run3 = p3.add_run()
    run3.text = "Powered by Anthropic Claude"
    run3.font.size = Pt(13)
    run3.font.color.rgb = ORANGE
    run3.font.name = "Calibri"
    run3.font.italic = True

    # Decorative teal dots (small circles as accent)
    for i, xpos in enumerate([Inches(1.5), Inches(2.0), Inches(2.5)]):
        dot = slide.shapes.add_shape(9, xpos, Inches(6.5), Inches(0.12), Inches(0.12))
        dot.fill.solid()
        dot.fill.fore_color.rgb = TEAL if i == 1 else ORANGE
        dot.line.fill.background()


def slide_02(prs):
    """What Is Claude Code?"""
    slide = blank_slide(prs)
    add_background(slide)
    add_title(slide, "What Is Claude Code?")
    add_orange_rule(slide)

    bullets = [
        "Anthropic's agentic CLI — runs Claude directly in terminal, IDE, or browser",
        "Not just chat — an autonomous coding agent that reads, writes, runs, and searches your codebase",
        "Executes shell commands, tests, and build pipelines",
        "Spawns sub-agents to parallelise complex work",
        "Remembers context across sessions via memory files",
    ]
    add_bullet_box(slide, bullets, MARGIN, CONTENT_Y, W - 2 * MARGIN, Inches(2.5), font_size=15)

    # "Available on" line
    avail_tb = slide.shapes.add_textbox(MARGIN, Inches(3.85), W - 2 * MARGIN, Inches(0.45))
    tf = avail_tb.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Available on:  "
    run.font.size = Pt(13)
    run.font.color.rgb = GREY
    run.font.name = "Calibri"
    run2 = p.add_run()
    run2.text = "CLI · Desktop App · claude.ai/code · VS Code · JetBrains"
    run2.font.size = Pt(13)
    run2.font.color.rgb = TEAL
    run2.font.name = "Calibri"
    run2.font.bold = True

    # Key mental model highlight box
    add_highlight_box(
        slide,
        "Key mental model:  \"Claude Code is a senior engineer at the keyboard — "
        "not a text generator. You describe intent; it executes.\"",
        MARGIN, Inches(4.5), W - 2 * MARGIN, Inches(1.1),
        font_size=14
    )


def slide_03(prs):
    """Core Features."""
    slide = blank_slide(prs)
    add_background(slide)
    add_title(slide, "Core Features at a Glance")
    add_orange_rule(slide)

    headers = ["Feature", "What It Does"]
    rows = [
        ["Agentic file editing",  "Reads, diffs, and edits files with surgical precision"],
        ["Bash execution",        "Runs tests, installs packages, starts servers"],
        ["Web search & fetch",    "Researches docs and APIs in real time"],
        ["Plan Mode",             "Designs architecture before writing any code"],
        ["Sub-agents",            "Spawns specialised agents for parallel work"],
        ["Skills (/commands)",    "Reusable slash-command scripts you invoke by name"],
        ["MCP servers",           "Plug-in integrations: databases, Slack, Jira, browsers"],
        ["Memory system",         "Persists facts, preferences, decisions across sessions"],
        ["Hooks",                 "Shell scripts that fire on tool events"],
        ["CLAUDE.md",             "Project-level instructions loaded into every session"],
    ]
    col_widths = [Inches(3.2), Inches(9.0)]
    add_table(slide, headers, rows,
              MARGIN, CONTENT_Y,
              W - 2 * MARGIN, Inches(5.6),
              col_widths=col_widths)


def slide_04(prs):
    """Plan Mode."""
    slide = blank_slide(prs)
    add_background(slide)
    add_title(slide, "Plan Mode — Design Before You Build")
    add_orange_rule(slide)

    col_w = (W - 2 * MARGIN - Inches(0.3)) / 2

    # Left column header
    add_textbox(slide, "What It Does", MARGIN, CONTENT_Y, col_w, Inches(0.4),
                font_size=14, bold=True, color=ORANGE)

    left_bullets = [
        "Read-only architect mode  (Shift+Tab+Tab  or  /plan)",
        "Explores the codebase, asks clarifying questions",
        "Produces a structured plan — without touching a single file",
        "Plan saved to  .claude/plans/  for your review",
    ]
    add_bullet_box(slide, left_bullets, MARGIN, Inches(1.75), col_w, Inches(2.0), font_size=13)

    # Right column header
    rx = MARGIN + col_w + Inches(0.3)
    add_textbox(slide, "Why Use It", rx, CONTENT_Y, col_w, Inches(0.4),
                font_size=14, bold=True, color=ORANGE)

    right_bullets = [
        "Prevents expensive wrong-direction implementations",
        "Forces alignment before work starts",
        "Generates a reviewable, editable plan file",
    ]
    add_bullet_box(slide, right_bullets, rx, Inches(1.75), col_w, Inches(1.5), font_size=13)

    # Code block
    code = """/plan

I want to build a robotics news aggregator page.
The backend has an existing news pattern to replicate.
No AWS credentials — use the direct Anthropic API.
Consider security at every layer. Use the
backend-architect and security-auditor agents."""
    add_code_block(slide, code, MARGIN, Inches(3.9), W - 2 * MARGIN, Inches(2.85), font_size=11)


def slide_05(prs):
    """Sub-Agents."""
    slide = blank_slide(prs)
    add_background(slide)
    add_title(slide, "Sub-Agents — Parallelism & Specialisation")
    add_orange_rule(slide)

    # Intro line
    add_textbox(slide, "Claude Code can spawn child agents, each with a specialised role.",
                MARGIN, CONTENT_Y, W - 2 * MARGIN, Inches(0.4),
                font_size=14, color=GREY)

    headers = ["Agent", "Best For"]
    rows = [
        ["Explore",              "Fast codebase mapping, file search"],
        ["Plan",                 "Architecture design, implementation strategy"],
        ["fullstack-developer",  "Features spanning DB + API + frontend"],
        ["code-reviewer",        "Quality, security, SOLID principles"],
        ["debugger",             "Root cause analysis, stack trace investigation"],
        ["python-pro",           "Type-safe async Python, FastAPI, pytest"],
        ["devops-engineer",      "CI/CD, Docker, cloud infrastructure"],
        ["general-purpose",      "Multi-step research, open-ended tasks"],
    ]
    col_widths = [Inches(3.5), Inches(8.7)]
    add_table(slide, headers, rows,
              MARGIN, Inches(1.72),
              W - 2 * MARGIN, Inches(4.4),
              col_widths=col_widths)

    # Key principle
    add_highlight_box(
        slide,
        "Key principle:  Independent tasks run in parallel. Dependent tasks run sequentially.",
        MARGIN, Inches(6.3), W - 2 * MARGIN, Inches(0.75),
        font_size=13
    )


def slide_06(prs):
    """Skills."""
    slide = blank_slide(prs)
    add_background(slide)
    add_title(slide, "Skills — Reusable Slash Commands")
    add_orange_rule(slide)

    col_w = (W - 2 * MARGIN - Inches(0.3)) / 2

    # Left: What are Skills?
    add_textbox(slide, "What are Skills?", MARGIN, CONTENT_Y, col_w, Inches(0.4),
                font_size=14, bold=True, color=ORANGE)

    left_bullets = [
        "Markdown files in  .claude/commands/",
        "Define a custom /command you invoke by name",
        "Can reference conversation history, files, and tools",
        "Shareable across your team",
    ]
    add_bullet_box(slide, left_bullets, MARGIN, Inches(1.75), col_w, Inches(1.8), font_size=13)

    # Right: Example skill
    rx = MARGIN + col_w + Inches(0.3)
    add_textbox(slide, "Example skill file:", rx, CONTENT_Y, col_w, Inches(0.4),
                font_size=14, bold=True, color=ORANGE)

    code = """---
description: Summarise session → append to CLAUDE.md
---

Review this conversation. Extract:
1. Architectural decisions made
2. Patterns established
3. Files created or changed
4. User preferences

Append a dated summary block to CLAUDE.md."""
    add_code_block(slide, code, rx, Inches(1.75), col_w, Inches(3.8), font_size=10)

    # Bottom invoke line
    invoke_tb = slide.shapes.add_textbox(MARGIN, Inches(6.45), W - 2 * MARGIN, Inches(0.55))
    tf = invoke_tb.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Invoke with:  "
    run.font.size = Pt(14)
    run.font.color.rgb = GREY
    run.font.name = "Calibri"
    run2 = p.add_run()
    run2.text = "/save-to-claude-md"
    run2.font.size = Pt(14)
    run2.font.color.rgb = TEAL
    run2.font.name = "Courier New"
    run2.font.bold = True


def slide_07(prs):
    """MCP Servers."""
    slide = blank_slide(prs)
    add_background(slide)
    add_title(slide, "MCP — Plug-In Integrations for Claude")
    add_orange_rule(slide)

    add_textbox(slide, "Model Context Protocol — an open standard for tool plugins",
                MARGIN, CONTENT_Y, W - 2 * MARGIN, Inches(0.35),
                font_size=13, color=TEAL)

    # Install code block
    code = """claude mcp add github    -- npx @anthropic/mcp-server-github
claude mcp add postgres  -- npx @anthropic/mcp-server-postgres"""
    add_code_block(slide, code, MARGIN, Inches(1.7), W - 2 * MARGIN, Inches(0.85), font_size=11)

    headers = ["MCP", "What Claude Can Do"]
    rows = [
        ["GitHub",               "Open PRs, comment on issues, create branches"],
        ["PostgreSQL",           "Query and migrate databases directly"],
        ["Slack",                "Post messages, read channels"],
        ["Browser (Puppeteer)",  "Navigate pages, screenshots, fill forms"],
        ["Jira / Linear",        "Create and update tickets"],
        ["Filesystem",           "Extended file operations"],
    ]
    col_widths = [Inches(2.8), Inches(9.4)]
    add_table(slide, headers, rows,
              MARGIN, Inches(2.7),
              W - 2 * MARGIN, Inches(3.3),
              col_widths=col_widths)

    add_highlight_box(
        slide,
        "MCP turns Claude from a code editor into a full workflow automation engine.",
        MARGIN, Inches(6.15), W - 2 * MARGIN, Inches(0.75),
        font_size=14
    )


def slide_08(prs):
    """app.aitmpl.com."""
    slide = blank_slide(prs)
    add_background(slide)
    add_title(slide, "app.aitmpl.com — The Claude Code Marketplace")
    add_orange_rule(slide)

    add_textbox(slide,
                "Community-built Skills, Agents, MCPs, Hooks, Commands, and Settings",
                MARGIN, CONTENT_Y, W - 2 * MARGIN, Inches(0.4),
                font_size=13, color=TEAL)

    col_w = (W - 2 * MARGIN - Inches(0.4)) / 2

    # Left: What you can browse
    add_textbox(slide, "What You Can Browse & Install", MARGIN, Inches(1.75), col_w, Inches(0.4),
                font_size=13, bold=True, color=ORANGE)
    browse_bullets = [
        "Skills — slash commands for common workflows",
        "Agents — specialised sub-agent definitions",
        "MCP configs — ready-to-use server configurations",
        "Hooks — event-triggered shell scripts",
        "Settings — best-practice presets",
    ]
    add_bullet_box(slide, browse_bullets, MARGIN, Inches(2.22), col_w, Inches(2.2), font_size=13)

    # Right: Featured integrations
    rx = MARGIN + col_w + Inches(0.4)
    add_textbox(slide, "Featured Integrations", rx, Inches(1.75), col_w, Inches(0.4),
                font_size=13, bold=True, color=ORANGE)
    feat_bullets = [
        "BrainGrid — AI product planner → generates scoped tasks Claude can execute",
        "Bright Data — Full web scraping stack (perfect for news aggregators)",
        "TinyFish — AI web agent platform",
        "ClaudeKit — Pre-built agents and skills toolkit",
    ]
    add_bullet_box(slide, feat_bullets, rx, Inches(2.22), col_w, Inches(2.2), font_size=13)

    # How to use
    add_highlight_box(
        slide,
        "How to use:  Browse → Add to Stack → Download → "
        "Copy to  .claude/commands/  or run  claude mcp add",
        MARGIN, Inches(4.6), W - 2 * MARGIN, Inches(0.85),
        font_size=13
    )

    # URL accent
    url_tb = slide.shapes.add_textbox(MARGIN, Inches(5.65), W - 2 * MARGIN, Inches(0.5))
    tf = url_tb.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "app.aitmpl.com"
    run.font.size = Pt(22)
    run.font.bold = True
    run.font.color.rgb = TEAL
    run.font.name = "Calibri"


def slide_09(prs):
    """CLAUDE.md."""
    slide = blank_slide(prs)
    add_background(slide)
    add_title(slide, "CLAUDE.md — Your Project's Persistent Brain")
    add_orange_rule(slide)

    add_textbox(slide, "Automatically injected into every Claude Code session",
                MARGIN, CONTENT_Y, W - 2 * MARGIN, Inches(0.35),
                font_size=13, color=TEAL)

    code = """# Project: Tech News Aggregator

## Stack
- Backend: FastAPI, Python 3.11, SQLite, Anthropic SDK
- Frontend: React 19, TypeScript, Vite, TanStack Query

## Conventions
- No comments unless the WHY is non-obvious
- Parameterised queries only — never string-format SQL
- Allowlist-validate all external inputs

## Architecture Decisions
- Dual-provider: settings.use_bedrock gates Bedrock vs Anthropic
- env_ignore_empty=True (OS empty vars don't override .env)

## Do Not
- Remove module-level Anthropic import (needed for test mocking)"""
    add_code_block(slide, code, MARGIN, Inches(1.7), W - 2 * MARGIN, Inches(4.7), font_size=10.5)

    add_highlight_box(
        slide,
        "Rule of thumb:  If you'd have to re-explain it next session, it belongs in CLAUDE.md.",
        MARGIN, Inches(6.6), W - 2 * MARGIN, Inches(0.65),
        font_size=13
    )


def slide_10(prs):
    """Workshop Project Overview."""
    slide = blank_slide(prs)
    add_background(slide)
    add_title(slide, "What We Built: Tech News Aggregator")
    add_orange_rule(slide)

    col_w = (W - 2 * MARGIN - Inches(0.35)) / 2

    # Left: Architecture
    add_textbox(slide, "Architecture", MARGIN, CONTENT_Y, col_w, Inches(0.4),
                font_size=14, bold=True, color=ORANGE)

    arch_code = """backend/
  services/agent_service.py   ← AI agent (dual provider)
  tools/news_tools.py         ← RSS fetch + tool schemas
  tools/robotics_tools.py     ← Robotics RSS feeds
  api/endpoints/robotics.py   ← GET /robotics endpoint
  database/db.py              ← SQLite cache + migrations
frontend/
  pages/RoboticsPage.tsx      ← Robotics news page
  hooks/useRoboticsNews.ts    ← TanStack Query hook
  App.tsx                     ← View toggle"""
    add_code_block(slide, arch_code, MARGIN, Inches(1.75), col_w, Inches(4.9), font_size=9.5)

    # Right: Capabilities
    rx = MARGIN + col_w + Inches(0.35)
    add_textbox(slide, "Capabilities", rx, CONTENT_Y, col_w, Inches(0.4),
                font_size=14, bold=True, color=ORANGE)

    cap_bullets = [
        "Chat interface → Claude agent → RSS feeds → SQLite cache",
        "Dual AI provider: AWS Bedrock OR Anthropic API key",
        "Robotics page with 5 sub-topic filters",
        "Streaming responses via SSE",
        "53 passing tests",
    ]
    add_bullet_box(slide, cap_bullets, rx, Inches(1.75), col_w, Inches(3.0), font_size=13)

    # Teal accent badge for "53 tests"
    badge = slide.shapes.add_shape(5, rx + col_w - Inches(2.1), Inches(4.9), Inches(2.0), Inches(0.65))
    badge.fill.solid()
    badge.fill.fore_color.rgb = RGBColor(0x0E, 0x3D, 0x36)
    badge.line.color.rgb = TEAL
    badge.line.width = Pt(1.5)
    tf = badge.text_frame
    tf.margin_left = Inches(0.1)
    tf.margin_top = Inches(0.08)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "53 Tests Passing"
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = TEAL
    run.font.name = "Calibri"


def slide_11(prs):
    """Step-by-Step."""
    slide = blank_slide(prs)
    add_background(slide)
    add_title(slide, "Building With Claude Code — Step by Step")
    add_orange_rule(slide)

    steps = [
        ("1. EXPLORE",    '"Familiarise yourself with the codebase. Why is AWS used?"'),
        ("2. PLAN",       '/plan  →  Request robotics feature + AWS bypass + security review'),
        ("3. REVIEW",     'Read .claude/plans/your-plan.md → edit → approve'),
        ("4. IMPLEMENT",  'Agent executes Steps 1–10 (backend), then 11–16 (frontend)'),
        ("5. TEST",       '"Run the full test suite and fix any failures"'),
        ("6. DEBUG",      '"The robotics page shows no articles found" → Claude traces + fixes'),
        ("7. SAVE",       '/save-to-claude-md → decisions appended to CLAUDE.md'),
    ]

    card_w = (W - 2 * MARGIN - Inches(0.15) * 6) / 7
    card_h = Inches(4.7)
    card_y = CONTENT_Y

    colors = [ORANGE, TEAL, ORANGE, TEAL, ORANGE, TEAL, ORANGE]

    for i, (step_name, desc) in enumerate(steps):
        cx = MARGIN + i * (card_w + Inches(0.15))

        # Card background
        card = slide.shapes.add_shape(5, cx, card_y, card_w, card_h)
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(0x13, 0x1C, 0x2A)
        card.line.color.rgb = colors[i]
        card.line.width = Pt(1.5)

        # Step number + name
        hdr = slide.shapes.add_textbox(cx + Inches(0.05), card_y + Inches(0.1),
                                        card_w - Inches(0.1), Inches(0.75))
        tf = hdr.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = step_name
        run.font.size = Pt(10)
        run.font.bold = True
        run.font.color.rgb = colors[i]
        run.font.name = "Calibri"

        # Description
        dtb = slide.shapes.add_textbox(cx + Inches(0.06), card_y + Inches(0.9),
                                        card_w - Inches(0.12), card_h - Inches(1.0))
        tf2 = dtb.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.alignment = PP_ALIGN.CENTER
        run2 = p2.add_run()
        run2.text = desc
        run2.font.size = Pt(9)
        run2.font.color.rgb = GREY
        run2.font.name = "Calibri"


def slide_12(prs):
    """Useful Agents."""
    slide = blank_slide(prs)
    add_background(slide)
    add_title(slide, "Best Agents for This Project")
    add_orange_rule(slide)

    headers = ["Agent", "Used For", "Example Prompt"]
    rows = [
        ["Plan",               "Architecture design",           "\"Design dual-provider agent service with security review\""],
        ["python-pro",         "Type-safe FastAPI code",        "\"Add async type hints and pytest coverage to db.py\""],
        ["code-reviewer",      "Pre-merge quality gate",        "\"Review agent_service.py for security and SOLID compliance\""],
        ["debugger",           "Root cause analysis",           "\"Robotics page returns no articles — trace why\""],
        ["Explore",            "Codebase mapping",              "\"How does the existing news endpoint work?\""],
        ["fullstack-developer","Cross-layer features",          "\"Build the complete Robotics feature DB → API → UI\""],
    ]
    col_widths = [Inches(2.2), Inches(2.8), Inches(7.2)]
    add_table(slide, headers, rows,
              MARGIN, CONTENT_Y,
              W - 2 * MARGIN, Inches(4.4),
              col_widths=col_widths)

    add_highlight_box(
        slide,
        "From aitmpl.com:  look for a web-scraping agent, a security auditor skill, and a test generator skill",
        MARGIN, Inches(5.95), W - 2 * MARGIN, Inches(0.8),
        font_size=13
    )


def slide_13(prs):
    """Best Practices."""
    slide = blank_slide(prs)
    add_background(slide)
    add_title(slide, "Best Practices Summary")
    add_orange_rule(slide)

    practices = [
        ("Plan First",           "Always enter Plan Mode before any implementation"),
        ("Write CLAUDE.md Early","Define stack, conventions, and \"do not\" rules upfront"),
        ("Parallelise with Agents","Code review + implementation can run simultaneously"),
        ("Security by Default",  "Allowlist inputs · Parameterise queries · Never log secrets"),
        ("Test Before Done",     "Run the full suite and fix failures before declaring complete"),
        ("Save Every Session",   "/save-to-claude-md keeps your project memory current"),
    ]

    card_w = (W - 2 * MARGIN - Inches(0.25)) / 3
    card_h = Inches(2.3)
    acc_colors = [ORANGE, TEAL, ORANGE, TEAL, ORANGE, TEAL]

    for i, (title, desc) in enumerate(practices):
        row = i // 3
        col = i % 3
        cx = MARGIN + col * (card_w + Inches(0.125))
        cy = CONTENT_Y + row * (card_h + Inches(0.2))

        # Card
        card = slide.shapes.add_shape(5, cx, cy, card_w, card_h)
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(0x13, 0x1C, 0x2A)
        card.line.color.rgb = acc_colors[i]
        card.line.width = Pt(1.5)

        # Top accent bar
        bar = slide.shapes.add_shape(1, cx, cy, card_w, Inches(0.06))
        bar.fill.solid()
        bar.fill.fore_color.rgb = acc_colors[i]
        bar.line.fill.background()

        # Title
        t_tb = slide.shapes.add_textbox(cx + Inches(0.12), cy + Inches(0.14),
                                         card_w - Inches(0.24), Inches(0.5))
        tf = t_tb.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = title
        run.font.size = Pt(13)
        run.font.bold = True
        run.font.color.rgb = acc_colors[i]
        run.font.name = "Calibri"

        # Description
        d_tb = slide.shapes.add_textbox(cx + Inches(0.12), cy + Inches(0.7),
                                         card_w - Inches(0.24), card_h - Inches(0.85))
        tf2 = d_tb.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        run2 = p2.add_run()
        run2.text = desc
        run2.font.size = Pt(12)
        run2.font.color.rgb = GREY
        run2.font.name = "Calibri"


def slide_14(prs):
    """Hands-On Checklist."""
    slide = blank_slide(prs)
    add_background(slide)
    add_title(slide, "Workshop Hands-On Checklist")
    add_orange_rule(slide)

    add_textbox(slide, "Follow these steps during the session",
                MARGIN, CONTENT_Y, W - 2 * MARGIN, Inches(0.35),
                font_size=13, color=TEAL)

    checklist = [
        "Clone the repo and open it with claude in the terminal",
        "Read CLAUDE.md — understand what's already documented",
        "Run ./start.sh — confirm backend + frontend are running",
        "Run pytest — confirm 53 tests pass",
        "Enter Plan Mode — ask for a new feature (e.g. \"add a cybersecurity news tab\")",
        "Review the generated plan file in .claude/plans/",
        "Approve and let the agent implement",
        "Watch sub-agents spawn in the terminal output",
        "Run tests again — fix failures with Claude's help",
        "Browse app.aitmpl.com — find one skill or agent to install",
        "Install it and invoke it  (/your-new-skill)",
        "Run /save-to-claude-md — verify CLAUDE.md is updated",
        "Open http://localhost:5173 — test the feature you built",
    ]

    col_w = (W - 2 * MARGIN - Inches(0.3)) / 2
    half = len(checklist) // 2 + len(checklist) % 2

    for col_idx in range(2):
        cx = MARGIN + col_idx * (col_w + Inches(0.3))
        items = checklist[col_idx * half: col_idx * half + half]
        y_offset = Inches(1.72)
        for item in items:
            item_tb = slide.shapes.add_textbox(cx, y_offset, col_w, Inches(0.42))
            tf = item_tb.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            # Checkbox char
            run1 = p.add_run()
            run1.text = "\u25a1  "
            run1.font.size = Pt(12)
            run1.font.color.rgb = ORANGE
            run1.font.name = "Calibri"
            run2 = p.add_run()
            run2.text = item
            run2.font.size = Pt(12)
            run2.font.color.rgb = GREY
            run2.font.name = "Calibri"
            y_offset += Inches(0.44)


def slide_15(prs):
    """Key Takeaways (closing slide)."""
    slide = blank_slide(prs)
    add_background(slide)
    add_title(slide, "Key Takeaways", font_size=36)
    add_orange_rule(slide)

    # Large central quote
    quote_box = slide.shapes.add_shape(1,
        MARGIN, Inches(1.35), W - 2 * MARGIN, Inches(0.95))
    quote_box.fill.solid()
    quote_box.fill.fore_color.rgb = RGBColor(0x1A, 0x10, 0x07)
    quote_box.line.color.rgb = ORANGE
    quote_box.line.width = Pt(1.5)

    tf = quote_box.text_frame
    tf.margin_left = Inches(0.2)
    tf.margin_top = Inches(0.15)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = '"Claude Code is not autocomplete. It is a collaborator."'
    run.font.size = Pt(20)
    run.font.bold = True
    run.font.color.rgb = ORANGE
    run.font.name = "Calibri"
    run.font.italic = True

    # Six bullet takeaways
    takeaways = [
        "Plan before you code — Plan Mode is your architecture review, for free",
        "Agents specialise — the right agent beats one generic prompt",
        "Skills scale your workflow — 10 minutes of setup saves hours per week",
        "MCP connects everything — your codebase is no longer an island",
        "CLAUDE.md is your team wiki — write it as if the next developer is Claude",
        "app.aitmpl.com is your shortcut — the community has solved most patterns",
    ]

    col_w = (W - 2 * MARGIN - Inches(0.3)) / 2
    half = 3

    for col_idx in range(2):
        cx = MARGIN + col_idx * (col_w + Inches(0.3))
        items = takeaways[col_idx * half: col_idx * half + half]
        y_offset = Inches(2.5)
        for item in items:
            tb = slide.shapes.add_textbox(cx, y_offset, col_w, Inches(0.6))
            tf = tb.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            run1 = p.add_run()
            run1.text = "▸  "
            run1.font.size = Pt(13)
            run1.font.color.rgb = TEAL
            run1.font.name = "Calibri"
            run2 = p.add_run()
            run2.text = item
            run2.font.size = Pt(13)
            run2.font.color.rgb = GREY
            run2.font.name = "Calibri"
            y_offset += Inches(0.7)

    # Bottom "Thank you" bar
    bar = slide.shapes.add_shape(1, 0, H - Inches(0.65), W, Inches(0.65))
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(0x12, 0x1A, 0x24)
    bar.line.fill.background()

    footer_tb = slide.shapes.add_textbox(MARGIN, H - Inches(0.58), W - 2 * MARGIN, Inches(0.5))
    tf = footer_tb.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER

    run1 = p.add_run()
    run1.text = "Thank you!  |  "
    run1.font.size = Pt(13)
    run1.font.color.rgb = WHITE
    run1.font.name = "Calibri"

    run2 = p.add_run()
    run2.text = "github.com/anthropics/claude-code"
    run2.font.size = Pt(13)
    run2.font.color.rgb = TEAL
    run2.font.name = "Calibri"

    run3 = p.add_run()
    run3.text = "  |  "
    run3.font.size = Pt(13)
    run3.font.color.rgb = WHITE
    run3.font.name = "Calibri"

    run4 = p.add_run()
    run4.text = "docs.anthropic.com/claude-code"
    run4.font.size = Pt(13)
    run4.font.color.rgb = TEAL
    run4.font.name = "Calibri"


# ══════════════════════════════════════════════════════════════════════════════
#  MAIN
# ══════════════════════════════════════════════════════════════════════════════

def main():
    output_path = r"C:\Users\a.lacatusu\Desktop\Claude\ClaudeCodeLabCamp\Claude_Code_Workshop.pptx"

    prs = new_prs()

    builders = [
        slide_01, slide_02, slide_03, slide_04, slide_05,
        slide_06, slide_07, slide_08, slide_09, slide_10,
        slide_11, slide_12, slide_13, slide_14, slide_15,
    ]

    for i, builder in enumerate(builders, 1):
        print(f"  Building slide {i:02d}/15 — {builder.__name__} ...")
        builder(prs)

    prs.save(output_path)
    print(f"\nDone! Saved to {output_path}")


if __name__ == "__main__":
    main()
